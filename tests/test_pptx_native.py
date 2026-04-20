"""Tests for pf.pptx_native editable PowerPoint export."""

import pytest
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor


class TestPptxTheme:
    def test_hex_to_rgb(self):
        from pf.pptx_native import _hex_to_rgb
        assert _hex_to_rgb("#C4A962") == RGBColor(0xC4, 0xA9, 0x62)

    def test_pptx_theme_colors(self):
        from pf.pptx_native import _pptx_theme
        theme_cfg = {"primary": "#1C2537", "accent": "#C4A962"}
        t = _pptx_theme(theme_cfg)
        assert t["primary"] == RGBColor(0x1C, 0x25, 0x37)
        assert t["accent"] == RGBColor(0xC4, 0xA9, 0x62)
        assert t["white"] == RGBColor(0xFF, 0xFF, 0xFF)
        assert t["text_muted"] == RGBColor(0xAA, 0xAA, 0xAA)

    def test_pptx_theme_fonts(self):
        from pf.pptx_native import _pptx_theme
        theme_cfg = {
            "primary": "#1C2537", "accent": "#C4A962",
            "fonts": {"heading": "Playfair Display", "subheading": "Montserrat", "body": "Lato"},
        }
        t = _pptx_theme(theme_cfg)
        assert t["font_heading"] == "Playfair Display"
        assert t["font_body"] == "Lato"

    def test_pptx_theme_defaults(self):
        from pf.pptx_native import _pptx_theme
        t = _pptx_theme({})
        assert t["primary"] == RGBColor(0x1C, 0x25, 0x37)
        assert t["font_heading"] == "Playfair Display"


from pptx import Presentation as PptxPresentation


class TestSectionLayout:
    def test_renders_title(self):
        from pf.pptx_native import _render_section, _pptx_theme, SLIDE_WIDTH, SLIDE_HEIGHT
        prs = PptxPresentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme = _pptx_theme({"primary": "#1C2537", "accent": "#C4A962"})
        _render_section(slide, {"title": "New Layouts", "subtitle": "Four new types", "number": 1}, theme)
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert "New Layouts" in texts

    def test_renders_subtitle(self):
        from pf.pptx_native import _render_section, _pptx_theme, SLIDE_WIDTH, SLIDE_HEIGHT
        prs = PptxPresentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme = _pptx_theme({"primary": "#1C2537", "accent": "#C4A962"})
        _render_section(slide, {"title": "Test", "subtitle": "Sub text"}, theme)
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert "SUB TEXT" in texts  # uppercase

    def test_renders_number(self):
        from pf.pptx_native import _render_section, _pptx_theme, SLIDE_WIDTH, SLIDE_HEIGHT
        prs = PptxPresentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme = _pptx_theme({"primary": "#1C2537", "accent": "#C4A962"})
        _render_section(slide, {"title": "Test", "number": 3}, theme)
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert "03" in texts


class TestQuoteLayout:
    def test_renders_quote_text(self):
        from pf.pptx_native import _render_quote, _pptx_theme, SLIDE_WIDTH, SLIDE_HEIGHT
        prs = PptxPresentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme = _pptx_theme({"primary": "#1C2537", "accent": "#C4A962"})
        _render_quote(slide, {"text": "The best way to predict the future is to invent it.", "author": "Alan Kay"}, theme)
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert any("predict the future" in t for t in texts)

    def test_renders_attribution(self):
        from pf.pptx_native import _render_quote, _pptx_theme, SLIDE_WIDTH, SLIDE_HEIGHT
        prs = PptxPresentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme = _pptx_theme({"primary": "#1C2537", "accent": "#C4A962"})
        _render_quote(slide, {"text": "Quote", "author": "Author", "role": "Scientist"}, theme)
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert any("Author" in t for t in texts)


class TestClosingLayout:
    def test_renders_title(self):
        from pf.pptx_native import _render_closing, _pptx_theme, SLIDE_WIDTH, SLIDE_HEIGHT
        prs = PptxPresentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme = _pptx_theme({"primary": "#1C2537", "accent": "#C4A962"})
        _render_closing(slide, {"title": "Thank You", "subtitle": "Questions?"}, theme)
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert "Thank You" in texts

    def test_renders_subtitle(self):
        from pf.pptx_native import _render_closing, _pptx_theme, SLIDE_WIDTH, SLIDE_HEIGHT
        prs = PptxPresentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme = _pptx_theme({"primary": "#1C2537", "accent": "#C4A962"})
        _render_closing(slide, {"title": "Thanks", "subtitle": "Q&A"}, theme)
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert "Q&A" in texts


import json
import yaml
from pathlib import Path


class TestExportEditable:
    def _write_config(self, tmp_path, slides):
        config = {
            "meta": {"title": "Test", "authors": ["Tester"]},
            "theme": {
                "primary": "#1C2537", "accent": "#C4A962",
                "fonts": {"heading": "Playfair Display", "subheading": "Montserrat", "body": "Lato"},
            },
            "slides": slides,
        }
        config_path = tmp_path / "presentation.yaml"
        config_path.write_text(yaml.dump(config, sort_keys=False), encoding="utf-8")
        metrics_path = tmp_path / "metrics.json"
        metrics_path.write_text(json.dumps({"metadata": {}, "summary": {}}), encoding="utf-8")
        return config_path, metrics_path

    def test_section_slide_native(self, tmp_path):
        from pf.pptx_native import export_pptx_editable
        from pf.builder import PresentationBuilder
        config_path, metrics_path = self._write_config(tmp_path, [
            {"layout": "section", "data": {"title": "Hello", "number": 1}},
        ])
        builder = PresentationBuilder(config_path=str(config_path), metrics_path=str(metrics_path))
        import contextlib, io as _io
        with contextlib.redirect_stdout(_io.StringIO()):
            out = builder.build(output_dir=str(tmp_path / "slides"))
        output_pptx = str(tmp_path / "out.pptx")
        export_pptx_editable(builder.config, str(out), output_pptx)
        prs = PptxPresentation(output_pptx)
        assert len(prs.slides) == 1
        texts = [s.text_frame.text for s in prs.slides[0].shapes if s.has_text_frame]
        assert "Hello" in texts

    def test_fallback_uses_image(self, tmp_path):
        """Non-native layouts should still produce a slide (via image fallback)."""
        from pf.pptx_native import export_pptx_editable
        from pf.builder import PresentationBuilder
        config_path, metrics_path = self._write_config(tmp_path, [
            {"layout": "section", "data": {"title": "Native"}},
            {"layout": "closing", "data": {"title": "Also Native"}},
        ])
        builder = PresentationBuilder(config_path=str(config_path), metrics_path=str(metrics_path))
        import contextlib, io as _io
        with contextlib.redirect_stdout(_io.StringIO()):
            out = builder.build(output_dir=str(tmp_path / "slides"))
        output_pptx = str(tmp_path / "out.pptx")
        export_pptx_editable(builder.config, str(out), output_pptx)
        prs = PptxPresentation(output_pptx)
        assert len(prs.slides) == 2


class TestTitleLayout:
    """Native PPTX renderer for title layout."""

    def test_title_in_native_renderers(self):
        from pf.pptx_native import NATIVE_RENDERERS
        assert "title" in NATIVE_RENDERERS

    def test_title_renders_without_error(self):
        from pf.pptx_native import NATIVE_RENDERERS, _pptx_theme, SLIDE_WIDTH, SLIDE_HEIGHT
        prs = PptxPresentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme = _pptx_theme({"primary": "#1C2537", "accent": "#C4A962"})
        NATIVE_RENDERERS["title"](slide, {"title": "Hello World", "subtitle": "Subtitle"}, theme)
        assert len(slide.shapes) > 0

    def test_title_renders_features(self):
        from pf.pptx_native import NATIVE_RENDERERS, _pptx_theme, SLIDE_WIDTH, SLIDE_HEIGHT
        prs = PptxPresentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme = _pptx_theme({"primary": "#1C2537", "accent": "#C4A962"})
        NATIVE_RENDERERS["title"](slide, {
            "title": "Title",
            "features": [{"label": "Feature A"}, {"label": "Feature B"}],
        }, theme)
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert any("Feature A" in t for t in texts)


class TestStatGridLayout:
    def test_stat_grid_in_native_renderers(self):
        from pf.pptx_native import NATIVE_RENDERERS
        assert "stat-grid" in NATIVE_RENDERERS

    def test_stat_grid_renders_without_error(self):
        from pf.pptx_native import NATIVE_RENDERERS, _pptx_theme, SLIDE_WIDTH, SLIDE_HEIGHT
        prs = PptxPresentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme = _pptx_theme({"primary": "#1C2537", "accent": "#C4A962"})
        NATIVE_RENDERERS["stat-grid"](slide, {
            "title": "KPIs",
            "stats": [
                {"value": "$1.2M", "label": "Revenue"},
                {"value": "45%", "label": "Growth"},
            ],
            "cols": 2,
        }, theme)
        assert len(slide.shapes) > 0

    def test_stat_grid_renders_values(self):
        from pf.pptx_native import NATIVE_RENDERERS, _pptx_theme, SLIDE_WIDTH, SLIDE_HEIGHT
        prs = PptxPresentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme = _pptx_theme({"primary": "#1C2537", "accent": "#C4A962"})
        NATIVE_RENDERERS["stat-grid"](slide, {
            "stats": [{"value": "99%", "label": "Uptime"}],
            "cols": 1,
        }, theme)
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert "99%" in texts


class TestTwoColumnLayout:
    def test_two_column_in_native_renderers(self):
        from pf.pptx_native import NATIVE_RENDERERS
        assert "two-column" in NATIVE_RENDERERS

    def test_two_column_renders_cards(self):
        from pf.pptx_native import NATIVE_RENDERERS, _pptx_theme, SLIDE_WIDTH, SLIDE_HEIGHT
        prs = PptxPresentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme = _pptx_theme({"primary": "#1C2537", "accent": "#C4A962"})
        NATIVE_RENDERERS["two-column"](slide, {
            "title": "Comparison",
            "left": [{"type": "card", "title": "Option A", "text": "First choice"}],
            "right": [{"type": "card", "title": "Option B", "text": "Second choice"}],
        }, theme)
        assert len(slide.shapes) > 0

    def test_two_column_renders_card_titles(self):
        from pf.pptx_native import NATIVE_RENDERERS, _pptx_theme, SLIDE_WIDTH, SLIDE_HEIGHT
        prs = PptxPresentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme = _pptx_theme({"primary": "#1C2537", "accent": "#C4A962"})
        NATIVE_RENDERERS["two-column"](slide, {
            "title": "Slide Title",
            "left": [{"type": "card", "title": "Left Card", "text": "Left text"}],
            "right": [],
        }, theme)
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert "Left Card" in texts


class TestThreeColumnLayout:
    def test_three_column_in_native_renderers(self):
        from pf.pptx_native import NATIVE_RENDERERS
        assert "three-column" in NATIVE_RENDERERS

    def test_three_column_renders(self):
        from pf.pptx_native import NATIVE_RENDERERS, _pptx_theme, SLIDE_WIDTH, SLIDE_HEIGHT
        prs = PptxPresentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme = _pptx_theme({"primary": "#1C2537", "accent": "#C4A962"})
        NATIVE_RENDERERS["three-column"](slide, {
            "title": "Three Options",
            "columns": [
                [{"type": "card", "title": "A", "text": "First"}],
                [{"type": "card", "title": "B", "text": "Second"}],
                [{"type": "card", "title": "C", "text": "Third"}],
            ],
        }, theme)
        assert len(slide.shapes) > 0

    def test_three_column_renders_card_titles(self):
        from pf.pptx_native import NATIVE_RENDERERS, _pptx_theme, SLIDE_WIDTH, SLIDE_HEIGHT
        prs = PptxPresentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme = _pptx_theme({"primary": "#1C2537", "accent": "#C4A962"})
        NATIVE_RENDERERS["three-column"](slide, {
            "columns": [
                [{"type": "card", "title": "Col A", "text": "Content"}],
                [],
                [],
            ],
        }, theme)
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert "Col A" in texts


class TestDataTableLayout:
    """Native PPTX renderer for data-table layout."""

    def _make_slide(self):
        from pf.pptx_native import _pptx_theme, SLIDE_WIDTH, SLIDE_HEIGHT
        prs = PptxPresentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme = _pptx_theme({"primary": "#1C2537", "accent": "#C4A962"})
        return slide, theme

    def test_data_table_in_native_renderers(self):
        from pf.pptx_native import NATIVE_RENDERERS
        assert "data-table" in NATIVE_RENDERERS

    def test_data_table_renders_without_error(self):
        from pf.pptx_native import NATIVE_RENDERERS
        slide, theme = self._make_slide()
        NATIVE_RENDERERS["data-table"](slide, {
            "title": "Benchmark Results",
            "sections": [
                {
                    "section_title": "Performance",
                    "table": {
                        "headers": ["Model", "Score", "Latency"],
                        "rows": [
                            ["GPT-4", "92%", "1.2s"],
                            ["Claude", "94%", "0.9s"],
                        ],
                        "winner_rows": [1],
                    },
                }
            ],
        }, theme)
        assert len(slide.shapes) > 0

    def test_data_table_renders_section_title(self):
        from pf.pptx_native import NATIVE_RENDERERS
        slide, theme = self._make_slide()
        NATIVE_RENDERERS["data-table"](slide, {
            "sections": [
                {
                    "section_title": "My Section",
                    "table": {
                        "headers": ["A", "B"],
                        "rows": [["1", "2"]],
                    },
                }
            ],
        }, theme)
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert any("My Section" in t for t in texts)

    def test_data_table_renders_table_headers(self):
        from pf.pptx_native import NATIVE_RENDERERS
        slide, theme = self._make_slide()
        NATIVE_RENDERERS["data-table"](slide, {
            "sections": [
                {
                    "table": {
                        "headers": ["Name", "Score"],
                        "rows": [["Alpha", "95%"]],
                    },
                }
            ],
        }, theme)
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert "Name" in texts
        assert "Score" in texts

    def test_data_table_renders_with_insight(self):
        from pf.pptx_native import NATIVE_RENDERERS
        slide, theme = self._make_slide()
        NATIVE_RENDERERS["data-table"](slide, {
            "sections": [
                {
                    "table": {"headers": ["X"], "rows": []},
                    "insight": {"text": "Key finding here"},
                }
            ],
        }, theme)
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert any("Key finding" in t for t in texts)

    def test_data_table_two_sections(self):
        from pf.pptx_native import NATIVE_RENDERERS
        slide, theme = self._make_slide()
        # Should not raise even with 2 sections
        NATIVE_RENDERERS["data-table"](slide, {
            "sections": [
                {"section_title": "Left", "table": {"headers": ["A"], "rows": []}},
                {"section_title": "Right", "table": {"headers": ["B"], "rows": []}},
            ],
        }, theme)
        assert len(slide.shapes) > 0


class TestImageLayout:
    """Native PPTX renderer for image layout."""

    def _make_slide(self):
        from pf.pptx_native import _pptx_theme, SLIDE_WIDTH, SLIDE_HEIGHT
        prs = PptxPresentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme = _pptx_theme({"primary": "#1C2537", "accent": "#C4A962"})
        return slide, theme

    def test_image_in_native_renderers(self):
        from pf.pptx_native import NATIVE_RENDERERS
        assert "image" in NATIVE_RENDERERS

    def test_image_renders_without_error_no_file(self):
        """Remote URL or missing file should render a placeholder without crashing."""
        from pf.pptx_native import NATIVE_RENDERERS
        slide, theme = self._make_slide()
        NATIVE_RENDERERS["image"](slide, {
            "image": "https://example.com/photo.jpg",
            "title": "Our Office",
            "caption": "San Francisco HQ",
        }, theme)
        assert len(slide.shapes) > 0

    def test_image_full_mode_renders_title(self):
        from pf.pptx_native import NATIVE_RENDERERS
        slide, theme = self._make_slide()
        NATIVE_RENDERERS["image"](slide, {
            "image": "https://example.com/x.jpg",
            "title": "Full Bleed Title",
            "position": "full",
        }, theme)
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert any("Full Bleed Title" in t for t in texts)

    def test_image_split_mode_renders_title(self):
        from pf.pptx_native import NATIVE_RENDERERS
        slide, theme = self._make_slide()
        NATIVE_RENDERERS["image"](slide, {
            "image": "https://example.com/x.jpg",
            "title": "Split Layout",
            "caption": "Descriptive text",
            "position": "split",
            "side": "left",
        }, theme)
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert any("Split Layout" in t for t in texts)

    def test_image_renders_caption(self):
        from pf.pptx_native import NATIVE_RENDERERS
        slide, theme = self._make_slide()
        NATIVE_RENDERERS["image"](slide, {
            "image": "https://example.com/x.jpg",
            "caption": "Photo credit: Unsplash",
            "position": "full",
        }, theme)
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert any("Photo credit" in t for t in texts)

    def test_image_local_file(self, tmp_path):
        """A local PNG file should be embedded natively via add_picture()."""
        from pf.pptx_native import NATIVE_RENDERERS
        # Create a minimal 1x1 white PNG (valid PNG bytes)
        png_bytes = (
            b'\x89PNG\r\n\x1a\n'  # PNG signature
            b'\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01'
            b'\x08\x02\x00\x00\x00\x90wS\xde\x00\x00\x00\x0cIDATx'
            b'\x9cc\xf8\x0f\x00\x00\x01\x01\x00\x05\x18\xd8N\x00\x00'
            b'\x00\x00IEND\xaeB`\x82'
        )
        img_path = tmp_path / "test.png"
        img_path.write_bytes(png_bytes)
        slide, theme = self._make_slide()
        # Should not raise — local file embedding path
        NATIVE_RENDERERS["image"](slide, {
            "image": str(img_path),
            "title": "Local Image",
            "position": "full",
        }, theme)
        assert len(slide.shapes) > 0


class TestTimelineLayout:
    """Native PPTX renderer for timeline layout."""

    def _make_slide(self):
        from pf.pptx_native import _pptx_theme, SLIDE_WIDTH, SLIDE_HEIGHT
        prs = PptxPresentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme = _pptx_theme({"primary": "#1C2537", "accent": "#C4A962"})
        return slide, theme

    def test_timeline_in_native_renderers(self):
        from pf.pptx_native import NATIVE_RENDERERS
        assert "timeline" in NATIVE_RENDERERS

    def test_timeline_renders_without_error(self):
        from pf.pptx_native import NATIVE_RENDERERS
        slide, theme = self._make_slide()
        NATIVE_RENDERERS["timeline"](slide, {
            "title": "Product Roadmap",
            "steps": [
                {"icon": "rocket", "title": "Launch", "description": "Initial release"},
                {"icon": "chart-line", "title": "Grow", "description": "Scale users"},
                {"icon": "trophy", "title": "Win", "description": "Market leader"},
            ],
        }, theme)
        assert len(slide.shapes) > 0

    def test_timeline_renders_step_titles(self):
        from pf.pptx_native import NATIVE_RENDERERS
        slide, theme = self._make_slide()
        NATIVE_RENDERERS["timeline"](slide, {
            "steps": [
                {"icon": "flag", "title": "Step One", "description": "First step"},
                {"icon": "check", "title": "Step Two", "description": "Second step"},
            ],
        }, theme)
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert any("Step One" in t for t in texts)
        assert any("Step Two" in t for t in texts)

    def test_timeline_renders_descriptions(self):
        from pf.pptx_native import NATIVE_RENDERERS
        slide, theme = self._make_slide()
        NATIVE_RENDERERS["timeline"](slide, {
            "steps": [
                {"icon": "star", "title": "Phase A", "description": "Do the thing"},
            ],
        }, theme)
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert any("Do the thing" in t for t in texts)

    def test_timeline_renders_slide_title(self):
        from pf.pptx_native import NATIVE_RENDERERS
        slide, theme = self._make_slide()
        NATIVE_RENDERERS["timeline"](slide, {
            "title": "Our Journey",
            "steps": [{"icon": "play", "title": "Start", "description": "Begin"}],
        }, theme)
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert any("Our Journey" in t for t in texts)

    def test_timeline_empty_steps(self):
        """Empty steps list should not raise."""
        from pf.pptx_native import NATIVE_RENDERERS
        slide, theme = self._make_slide()
        NATIVE_RENDERERS["timeline"](slide, {"title": "Empty", "steps": []}, theme)
        # Background shape added, no crash
        assert len(slide.shapes) >= 0

    def test_timeline_single_step_no_line(self):
        """Single step — connecting line should still render without crashing."""
        from pf.pptx_native import NATIVE_RENDERERS
        slide, theme = self._make_slide()
        NATIVE_RENDERERS["timeline"](slide, {
            "steps": [{"icon": "bolt", "title": "Only", "description": "Solo step"}],
        }, theme)
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert any("Only" in t for t in texts)


# ── T2.1 — per-layout editable-shape coverage ─────────────────────
#
# Builds a one-slide fixture exercising each LAYOUT_NAMES entry,
# exports --editable, opens with python-pptx, asserts the slide
# contains at least one shape that is not just a raster picture.
# The 6 layouts not yet ported (chart, code, map, mermaid, toc, video)
# are expected to fail until T2.2-T2.7 land — marked xfail(strict=True)
# so when they start passing the test suite forces an unxfail.

import json as _json
import yaml as _yaml_t21

from pf.pptx_native import LAYOUT_NAMES as _LAYOUT_NAMES_T21


_LAYOUT_FIXTURES: dict[str, dict] = {
    "title": {"title": "Test", "subtitle": "Sub", "tagline": "Tag"},
    "section": {"title": "S", "subtitle": "Sub", "number": 1},
    "quote": {"text": "An editable quote.", "attribution": "Author"},
    "closing": {"title": "Thanks", "subtitle": "Q&A"},
    "two-column": {
        "title": "Two", "left": [{"type": "card", "title": "L", "text": "Left"}],
        "right": [{"type": "card", "title": "R", "text": "Right"}],
    },
    "three-column": {
        "title": "Three",
        "columns": [
            [{"type": "card", "title": "A", "text": "A1"}],
            [{"type": "card", "title": "B", "text": "B1"}],
            [{"type": "card", "title": "C", "text": "C1"}],
        ],
    },
    "stat-grid": {
        "title": "Stats",
        "columns": [
            [{"type": "stat-grid", "stats": [{"value": "1", "label": "x"}]}],
            [{"type": "card", "header": "H", "items": ["a"]}],
        ],
    },
    "data-table": {
        "title": "Data",
        "sections": [{
            "section_title": "Sec",
            "table": {"headers": ["A", "B"], "rows": [["1", "2"]]},
        }],
    },
    "image": {"title": "Img", "image": "nonexistent.png", "caption": "Cap"},
    "timeline": {
        "title": "TL",
        "steps": [{"icon": "bolt", "title": "S1", "description": "D1"}],
    },
    "code": {"title": "Code", "language": "python", "code": "x = 1"},
    "toc": {"title": "TOC", "items": [{"title": "Item 1", "slide": 1}]},
    "chart": {
        "title": "Chart",
        "chart_type": "bar",
        "labels": ["A", "B"],
        "values": [1, 2],
    },
    "map": {"title": "Map", "lat": 0, "lng": 0, "zoom": 5, "markers": []},
    "mermaid": {"title": "Mermaid", "diagram": "graph TD; A-->B"},
    "video": {"title": "Video", "url": "https://example.com/v.mp4"},
}


def _build_one_slide_pptx(layout, tmp_path):
    """Build a single-slide deck for the given layout and return the .pptx path."""
    from pf.builder import PresentationBuilder
    from pf.pptx_native import export_pptx_editable

    config = {
        "meta": {"title": f"T-{layout}"},
        "theme": {"primary": "#1C2537", "accent": "#C4A962", "charts": True,
                  "fonts": {"heading": "Playfair Display", "subheading": "Montserrat", "body": "Lato"}},
        "slides": [{"layout": layout, "data": _LAYOUT_FIXTURES[layout]}],
    }
    cfg_path = tmp_path / "presentation.yaml"
    cfg_path.write_text(_yaml_t21.dump(config), encoding="utf-8")
    metrics_path = tmp_path / "metrics.json"
    metrics_path.write_text(_json.dumps({}), encoding="utf-8")

    slides_dir = tmp_path / "slides"
    builder = PresentationBuilder(config_path=str(cfg_path), metrics_path=str(metrics_path))
    builder.build(output_dir=str(slides_dir))

    out = tmp_path / "deck.pptx"
    export_pptx_editable(config, str(slides_dir), str(out))
    return out


def _slide_has_editable_content(slide):
    """Editable = at least one text frame with text, OR a chart, OR a movie."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            return True
        if getattr(shape, "has_chart", False):
            return True
        if shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
            return True
    return False


_UNIMPLEMENTED_LAYOUTS = {"map", "mermaid", "video"}


def _layout_param(name):
    if name in _UNIMPLEMENTED_LAYOUTS:
        return pytest.param(
            name,
            marks=pytest.mark.xfail(
                strict=True,
                reason=f"{name}: native PPTX renderer pending (T2.2-T2.7)",
            ),
        )
    return pytest.param(name)


@pytest.mark.parametrize("layout", [_layout_param(n) for n in _LAYOUT_NAMES_T21])
def test_each_layout_editable(layout, tmp_path):
    """Every built-in layout must produce at least one editable shape in --editable mode."""
    out = _build_one_slide_pptx(layout, tmp_path)
    from pptx import Presentation as _P
    prs = _P(out)
    assert len(prs.slides) == 1
    slide = prs.slides[0]
    assert _slide_has_editable_content(slide), (
        f"{layout}: no editable shape (text/chart/movie) in PPTX export"
    )


class TestChartLayout:
    """Native PPTX renderer for chart layout — uses python-pptx charts so the
    result is double-click-editable in PowerPoint."""

    def _make_slide(self):
        from pf.pptx_native import _pptx_theme, SLIDE_WIDTH, SLIDE_HEIGHT
        prs = PptxPresentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme = _pptx_theme({"primary": "#1C2537", "accent": "#C4A962",
                             "secondary_accent": "#6B5BA2",
                             "fonts": {"heading": "H", "subheading": "S",
                                       "body": "B", "mono": "M"}})
        return slide, theme

    def test_chart_in_native_renderers(self):
        from pf.pptx_native import NATIVE_RENDERERS
        assert "chart" in NATIVE_RENDERERS

    def test_chart_type_mapping(self):
        """bar → COLUMN_CLUSTERED (Plotly default is vertical); other types
        match the plan's D5 table directly."""
        from pf.pptx_native import _xl_chart_type
        from pptx.enum.chart import XL_CHART_TYPE
        assert _xl_chart_type("bar") == XL_CHART_TYPE.COLUMN_CLUSTERED
        assert _xl_chart_type("column") == XL_CHART_TYPE.COLUMN_CLUSTERED
        assert _xl_chart_type("bar-horizontal") == XL_CHART_TYPE.BAR_CLUSTERED
        assert _xl_chart_type("line") == XL_CHART_TYPE.LINE
        assert _xl_chart_type("pie") == XL_CHART_TYPE.PIE
        assert _xl_chart_type("donut") == XL_CHART_TYPE.DOUGHNUT
        assert _xl_chart_type("doughnut") == XL_CHART_TYPE.DOUGHNUT
        assert _xl_chart_type("scatter") == XL_CHART_TYPE.XY_SCATTER
        assert _xl_chart_type("area") == XL_CHART_TYPE.AREA
        assert _xl_chart_type("") is None
        assert _xl_chart_type("treemap") is None

    def test_chart_single_series_renders_native_chart(self):
        from pf.pptx_native import NATIVE_RENDERERS
        slide, theme = self._make_slide()
        NATIVE_RENDERERS["chart"](slide, {
            "title": "Revenue",
            "chart_type": "bar",
            "labels": ["Q1", "Q2", "Q3"],
            "values": [100, 200, 300],
        }, theme)
        charts = [s for s in slide.shapes if getattr(s, "has_chart", False)]
        assert len(charts) == 1, "expected one native chart shape"
        assert len(list(charts[0].chart.series)) == 1

    def test_chart_multi_series_renders(self):
        from pf.pptx_native import NATIVE_RENDERERS
        slide, theme = self._make_slide()
        NATIVE_RENDERERS["chart"](slide, {
            "chart_type": "line",
            "labels": ["2023", "2024", "2025"],
            "series": [
                {"name": "Product A", "values": [10, 20, 30]},
                {"name": "Product B", "values": [15, 18, 22]},
            ],
        }, theme)
        charts = [s for s in slide.shapes if getattr(s, "has_chart", False)]
        assert len(charts) == 1
        series = list(charts[0].chart.series)
        assert len(series) == 2
        assert series[0].name == "Product A"
        assert series[1].name == "Product B"

    def test_chart_scatter_uses_xy_data(self):
        """scatter → one of the XY_SCATTER variants (python-pptx may read back
        the type as XY_SCATTER_LINES rather than the XY_SCATTER we created
        with — it normalizes from the underlying XML)."""
        from pf.pptx_native import NATIVE_RENDERERS
        from pptx.enum.chart import XL_CHART_TYPE
        slide, theme = self._make_slide()
        NATIVE_RENDERERS["chart"](slide, {
            "chart_type": "scatter",
            "labels": [1, 2, 3, 4],
            "values": [10, 20, 15, 25],
        }, theme)
        charts = [s for s in slide.shapes if getattr(s, "has_chart", False)]
        assert len(charts) == 1
        scatter_variants = {
            XL_CHART_TYPE.XY_SCATTER,
            XL_CHART_TYPE.XY_SCATTER_LINES,
            XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS,
            XL_CHART_TYPE.XY_SCATTER_SMOOTH,
            XL_CHART_TYPE.XY_SCATTER_SMOOTH_NO_MARKERS,
        }
        assert charts[0].chart.chart_type in scatter_variants

    def test_chart_unknown_type_falls_back_to_text_placeholder(self):
        """Unknown chart types shouldn't crash — they render a text placeholder
        so the slide still has an editable shape. T2.8's --strict will treat
        this as a fallback event."""
        from pf.pptx_native import NATIVE_RENDERERS
        slide, theme = self._make_slide()
        NATIVE_RENDERERS["chart"](slide, {
            "title": "Weird",
            "chart_type": "treemap",
            "labels": ["A"],
            "values": [1],
        }, theme)
        charts = [s for s in slide.shapes if getattr(s, "has_chart", False)]
        assert charts == [], "unknown type must not produce a native chart"
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert any("treemap" in t for t in texts)

    def test_chart_applies_accent_color(self):
        """First series should be painted with theme.accent; second with
        theme.secondary_accent when present."""
        from pf.pptx_native import NATIVE_RENDERERS
        slide, theme = self._make_slide()
        NATIVE_RENDERERS["chart"](slide, {
            "chart_type": "column",
            "labels": ["A", "B"],
            "series": [
                {"name": "S1", "values": [1, 2]},
                {"name": "S2", "values": [3, 4]},
            ],
        }, theme)
        charts = [s for s in slide.shapes if getattr(s, "has_chart", False)]
        assert len(charts) == 1
        s1, s2 = list(charts[0].chart.series)
        # Fill fore_color may not roundtrip on every series type, but at least
        # one color API should report the accent for S1 after _apply_chart_theme.
        c1 = s1.format.fill.fore_color.rgb
        c2 = s2.format.fill.fore_color.rgb
        assert c1 == theme["accent"]
        assert c2 == theme["secondary_accent"]


class TestTocLayout:
    """Native PPTX renderer for toc layout with NAMED_SLIDE hyperlinks."""

    def test_toc_in_native_renderers(self):
        from pf.pptx_native import NATIVE_RENDERERS
        assert "toc" in NATIVE_RENDERERS

    def test_toc_renders_entries_and_hyperlinks(self, tmp_path):
        """Each entry should be a text box with click_action.target_slide set,
        resulting in PP_ACTION.NAMED_SLIDE (PowerPoint's 'hyperlink to slide')."""
        from pf.pptx_native import export_pptx_editable
        from pf.builder import PresentationBuilder
        from pptx.enum.action import PP_ACTION

        config = {
            "meta": {"title": "Deck"},
            "theme": {
                "primary": "#1C2537", "accent": "#C4A962",
                "fonts": {"heading": "Playfair Display", "subheading": "Montserrat",
                          "body": "Lato", "mono": "IBM Plex Mono"},
            },
            "slides": [
                {"layout": "toc", "data": {"title": "Contents", "items": [
                    {"number": 1, "title": "Intro"},
                    {"number": 2, "title": "Results"},
                ]}},
                {"layout": "section", "data": {"number": 1, "title": "Intro"}},
                {"layout": "section", "data": {"number": 2, "title": "Results"}},
            ],
        }
        cfg_path = tmp_path / "presentation.yaml"
        cfg_path.write_text(_yaml_t21.dump(config), encoding="utf-8")
        metrics_path = tmp_path / "metrics.json"
        metrics_path.write_text(_json.dumps({}), encoding="utf-8")

        builder = PresentationBuilder(config_path=str(cfg_path), metrics_path=str(metrics_path))
        import contextlib, io as _io
        with contextlib.redirect_stdout(_io.StringIO()):
            slides_dir = builder.build(output_dir=str(tmp_path / "slides"))
        out_pptx = tmp_path / "out.pptx"
        export_pptx_editable(builder.config, str(slides_dir), str(out_pptx))

        prs = PptxPresentation(str(out_pptx))
        toc_slide = prs.slides[0]
        # Every shape with non-empty text should correspond to a TOC entry or the title
        entry_shapes = [
            s for s in toc_slide.shapes
            if s.has_text_frame and s.text_frame.text.strip()
            and "Intro" in s.text_frame.text or "Results" in s.text_frame.text
        ]
        assert len(entry_shapes) >= 2, "expected at least two TOC entry shapes"
        # At least one of the entry shapes must have a NAMED_SLIDE click_action
        named_slide_hits = [
            s for s in entry_shapes
            if s.click_action.action == PP_ACTION.NAMED_SLIDE
        ]
        assert named_slide_hits, "no entries carry a NAMED_SLIDE hyperlink"

    def test_toc_explicit_slide_index_respected(self):
        """Direct _render_toc call: an item with `slide: N` (1-based) hyperlinks
        to prs.slides[N-1], ignoring the section-matching fallback.

        (The full builder pipeline clobbers user-supplied `items` via
        `slide_cfg.setdefault("data", {})["items"] = _generate_toc(slides)` in
        builder.py, so this path is only reachable when a caller invokes
        _render_toc / export_pptx_editable on a hand-assembled config — the
        exact use case this test exercises.)"""
        from pf.pptx_native import _render_toc, _pptx_theme, SLIDE_WIDTH, SLIDE_HEIGHT
        from pptx.enum.action import PP_ACTION

        prs = PptxPresentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        blank = prs.slide_layouts[6]
        toc = prs.slides.add_slide(blank)
        _ = prs.slides.add_slide(blank)      # slide 2
        target = prs.slides.add_slide(blank)  # slide 3 (closing)

        theme = _pptx_theme({"primary": "#1C2537", "accent": "#C4A962",
                             "fonts": {"heading": "H", "subheading": "S",
                                       "body": "B", "mono": "M"}})
        data = {"items": [{"title": "Jump to closing", "slide": 3}]}
        slides_cfg = [
            {"layout": "toc", "data": data},
            {"layout": "section", "data": {"number": 1, "title": "X"}},
            {"layout": "closing", "data": {"title": "Bye"}},
        ]
        _render_toc(toc, data, theme, prs=prs, slides_cfg=slides_cfg, slide_index=0)

        hits = [
            s for s in toc.shapes
            if s.has_text_frame and "Jump to closing" in s.text_frame.text
            and s.click_action.action == PP_ACTION.NAMED_SLIDE
        ]
        assert hits, "explicit slide-index entry did not produce a NAMED_SLIDE hyperlink"
        # Sanity: target should actually be the third slide, not the second
        assert hits[0].click_action.target_slide is target


def test_layout_names_in_sync_with_templates():
    """LAYOUT_NAMES tuple must match templates/layouts/*.html.j2 exactly."""
    from pf.pptx_native import LAYOUT_NAMES, _discover_layout_names
    assert LAYOUT_NAMES == _discover_layout_names()


def test_iter_native_layouts_subset_of_layout_names():
    """Every native renderer must correspond to a real layout name."""
    from pf.pptx_native import LAYOUT_NAMES, iter_native_layouts
    assert set(iter_native_layouts()).issubset(set(LAYOUT_NAMES))

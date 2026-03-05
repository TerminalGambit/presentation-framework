"""
PowerPoint export for Presentation Framework.
Requires: pip install presentation-framework[pptx]

Renders each slide as a PNG screenshot via Playwright, then inserts
into a .pptx file using python-pptx. Output is pixel-perfect but
non-editable (image-based slides).
"""

import io
from pathlib import Path

try:
    from playwright.sync_api import sync_playwright
    PLAYWRIGHT_AVAILABLE = True
except ImportError:
    PLAYWRIGHT_AVAILABLE = False


def export_pptx(
    slides_dir: str,
    output_path: str,
    title: str = "Presentation",
):
    """Export built slides to a .pptx file.

    Each slide is rendered as a 1280x720 PNG screenshot and inserted
    as a full-bleed image into a PowerPoint slide.
    """
    if not PLAYWRIGHT_AVAILABLE:
        raise ImportError(
            "PPTX export requires Playwright. Install with:\n"
            "  pip install presentation-framework[pptx]\n"
            "  playwright install chromium"
        )

    from pptx import Presentation
    from pptx.util import Inches, Emu

    slides_path = Path(slides_dir).resolve()
    slide_files = sorted(slides_path.glob("slide_*.html"))

    if not slide_files:
        raise FileNotFoundError(f"No slide files found in {slides_dir}")

    # 1280x720 = 16:9 at 96 DPI
    SLIDE_WIDTH = Inches(13.333)   # 1280px at 96 DPI
    SLIDE_HEIGHT = Inches(7.5)     # 720px at 96 DPI

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]  # Blank layout

    with sync_playwright() as p:
        browser = p.chromium.launch()
        context = browser.new_context(
            viewport={"width": 1280, "height": 720},
        )

        for slide_file in slide_files:
            page = context.new_page()
            page.goto(f"file://{slide_file}")
            page.wait_for_load_state("networkidle")

            # Screenshot to bytes
            png_bytes = page.screenshot(full_page=False)
            page.close()

            # Add slide with full-bleed image
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                io.BytesIO(png_bytes),
                left=Emu(0),
                top=Emu(0),
                width=SLIDE_WIDTH,
                height=SLIDE_HEIGHT,
            )

        browser.close()

    prs.save(output_path)

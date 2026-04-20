"""
Editable PowerPoint export for Presentation Framework.
Requires: pip install presentation-framework[pptx]

Converts simple layouts (section, quote, closing) to native
python-pptx text boxes and shapes. Complex layouts fall back to
image-based rendering via Playwright.
"""

import inspect
import io
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ── Slide dimensions (16:9 at 96 DPI) ────────────────────────────
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


# ── Layout registry ──────────────────────────────────────────────
#
# Single source of truth for the 16 built-in layouts. Ordering matches
# `templates/layouts/*.html.j2` alphabetically; NATIVE_RENDERERS (below)
# is the editable-export dispatch table. The assert at import time
# catches drift between the tuple and the templates directory — a new
# layout file without a corresponding LAYOUT_NAMES entry is a bug.

LAYOUT_NAMES: tuple[str, ...] = (
    "chart",
    "closing",
    "code",
    "data-table",
    "image",
    "map",
    "mermaid",
    "quote",
    "section",
    "stat-grid",
    "three-column",
    "timeline",
    "title",
    "toc",
    "two-column",
    "video",
)


def _discover_layout_names() -> tuple[str, ...]:
    """Return the layout names discovered on disk under templates/layouts/."""
    layouts_dir = Path(__file__).resolve().parent.parent / "templates" / "layouts"
    if not layouts_dir.is_dir():
        return ()
    # Strip both .j2 and .html — file names are like 'chart.html.j2'
    return tuple(sorted(p.name.removesuffix(".html.j2") for p in layouts_dir.glob("*.html.j2")))


# Import-time guard: if someone adds/removes a layout template without
# updating LAYOUT_NAMES, crash loudly at import so the next test run or
# MCP tool call surfaces it immediately.
_discovered = _discover_layout_names()
if _discovered and _discovered != LAYOUT_NAMES:
    raise RuntimeError(
        "pptx_native.LAYOUT_NAMES is out of sync with templates/layouts/*.html.j2.\n"
        f"  expected: {_discovered}\n"
        f"  got:      {LAYOUT_NAMES}\n"
        "Update LAYOUT_NAMES (and, if adding a new layout, NATIVE_RENDERERS) "
        "in pf/pptx_native.py."
    )


def iter_native_layouts() -> tuple[str, ...]:
    """Return the layouts that currently have a NATIVE_RENDERERS entry.

    Layouts outside this tuple fall back to image rasterization in
    export_pptx_editable (and fail under --strict).
    """
    # NATIVE_RENDERERS is defined lower in this module; compute at call time.
    return tuple(sorted(NATIVE_RENDERERS.keys()))


# ── Theme conversion ─────────────────────────────────────────────

def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert '#RRGGBB' to python-pptx RGBColor."""
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _pptx_theme(theme_cfg: dict) -> dict:
    """Convert presentation.yaml theme to python-pptx primitives."""
    fonts = theme_cfg.get("fonts", {})
    result = {
        "primary": _hex_to_rgb(theme_cfg.get("primary", "#1C2537")),
        "accent": _hex_to_rgb(theme_cfg.get("accent", "#C4A962")),
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "text_muted": RGBColor(0xAA, 0xAA, 0xAA),
        "font_heading": fonts.get("heading", "Playfair Display"),
        "font_subheading": fonts.get("subheading", "Montserrat"),
        "font_body": fonts.get("body", "Lato"),
        "font_mono": fonts.get("mono", "IBM Plex Mono"),
    }
    if theme_cfg.get("secondary_accent"):
        result["secondary_accent"] = _hex_to_rgb(theme_cfg["secondary_accent"])
    return result


def _set_text(text_frame, text, font_name, font_size_pt, color, bold=False, alignment=PP_ALIGN.CENTER):
    """Set text on a shape's text_frame with styling."""
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size_pt)
    run.font.color.rgb = color
    run.font.bold = bold


def _add_bg(slide, color):
    """Set solid background color on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, color):
    """Add a colored rectangle shape."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()  # No border
    return shape


# ── Layout renderers ─────────────────────────────────────────────

def _render_section(slide, data: dict, theme: dict):
    """Render section divider: number + title + subtitle + accent bars."""
    _add_bg(slide, theme["primary"])

    center_x = SLIDE_WIDTH // 2
    y_cursor = Inches(2.0)

    # Top accent bar
    bar_w, bar_h = Inches(1.25), Inches(0.03)
    _add_rect(slide, center_x - bar_w // 2, y_cursor, bar_w, bar_h, theme["accent"])
    y_cursor += Inches(0.5)

    # Number (optional)
    if data.get("number") is not None:
        box_w, box_h = Inches(4), Inches(0.9)
        txBox = slide.shapes.add_textbox(center_x - box_w // 2, y_cursor, box_w, box_h)
        _set_text(txBox.text_frame, f"{data['number']:02d}", theme["font_heading"], 60, theme["accent"], bold=True)
        y_cursor += Inches(0.9)

    # Title
    box_w, box_h = Inches(10), Inches(1.2)
    txBox = slide.shapes.add_textbox(center_x - box_w // 2, y_cursor, box_w, box_h)
    _set_text(txBox.text_frame, data["title"], theme["font_heading"], 48, theme["white"], bold=True)
    y_cursor += Inches(1.2)

    # Subtitle (optional)
    if data.get("subtitle"):
        box_w, box_h = Inches(10), Inches(0.6)
        txBox = slide.shapes.add_textbox(center_x - box_w // 2, y_cursor, box_w, box_h)
        _set_text(txBox.text_frame, data["subtitle"].upper(), theme["font_subheading"], 18, theme["text_muted"])
        y_cursor += Inches(0.7)

    # Bottom accent bar
    _add_rect(slide, center_x - bar_w // 2, y_cursor, bar_w, bar_h, theme["accent"])


def _render_quote(slide, data: dict, theme: dict):
    """Render quote: quotation mark + text + divider + attribution."""
    _add_bg(slide, theme["primary"])
    center_x = SLIDE_WIDTH // 2
    y_cursor = Inches(1.5)

    # Quotation mark
    box_w, box_h = Inches(2), Inches(1.2)
    txBox = slide.shapes.add_textbox(center_x - box_w // 2, y_cursor, box_w, box_h)
    _set_text(txBox.text_frame, "\u201C", theme["font_heading"], 96, theme["accent"])
    y_cursor += Inches(1.0)

    # Quote text
    box_w, box_h = Inches(9), Inches(2.0)
    txBox = slide.shapes.add_textbox(center_x - box_w // 2, y_cursor, box_w, box_h)
    _set_text(txBox.text_frame, data["text"], theme["font_heading"], 28, theme["white"])
    txBox.text_frame.word_wrap = True
    y_cursor += Inches(1.8)

    # Divider
    bar_w, bar_h = Inches(0.8), Inches(0.03)
    _add_rect(slide, center_x - bar_w // 2, y_cursor, bar_w, bar_h, theme["accent"])
    y_cursor += Inches(0.4)

    # Attribution
    parts = []
    if data.get("author"):
        parts.append(data["author"])
    if data.get("role"):
        parts.append(data["role"])
    if parts:
        box_w, box_h = Inches(8), Inches(0.8)
        txBox = slide.shapes.add_textbox(center_x - box_w // 2, y_cursor, box_w, box_h)
        _set_text(txBox.text_frame, " — ".join(parts), theme["font_subheading"], 16, theme["text_muted"])


def _render_closing(slide, data: dict, theme: dict):
    """Render closing: title + subtitle + divider."""
    _add_bg(slide, theme["primary"])
    center_x = SLIDE_WIDTH // 2
    y_cursor = Inches(2.2)

    # Title
    box_w, box_h = Inches(10), Inches(1.5)
    txBox = slide.shapes.add_textbox(center_x - box_w // 2, y_cursor, box_w, box_h)
    _set_text(txBox.text_frame, data["title"], theme["font_heading"], 60, theme["accent"], bold=True)
    y_cursor += Inches(1.5)

    # Divider
    bar_w, bar_h = Inches(1.5), Inches(0.03)
    _add_rect(slide, center_x - bar_w // 2, y_cursor, bar_w, bar_h, theme["accent"])
    y_cursor += Inches(0.5)

    # Subtitle
    if data.get("subtitle"):
        box_w, box_h = Inches(10), Inches(0.8)
        txBox = slide.shapes.add_textbox(center_x - box_w // 2, y_cursor, box_w, box_h)
        _set_text(txBox.text_frame, data["subtitle"], theme["font_subheading"], 20, theme["text_muted"])


def _render_title(slide, data: dict, theme: dict):
    """Render title slide: hero title + subtitle + optional feature labels."""
    _add_bg(slide, theme["primary"])
    center_x = SLIDE_WIDTH // 2

    # Hero title
    box_w, box_h = Inches(10), Inches(2)
    txBox = slide.shapes.add_textbox(center_x - box_w // 2, Inches(1.8), box_w, box_h)
    _set_text(txBox.text_frame, data.get("title", ""), theme["font_heading"], 60, theme["accent"], bold=True)
    txBox.text_frame.word_wrap = True

    # Subtitle
    if data.get("subtitle"):
        box_w, box_h = Inches(10), Inches(0.8)
        txBox = slide.shapes.add_textbox(center_x - box_w // 2, Inches(3.9), box_w, box_h)
        _set_text(txBox.text_frame, data["subtitle"], theme["font_subheading"], 20, theme["text_muted"])
        txBox.text_frame.word_wrap = True

    # Features (icon labels across bottom)
    features = data.get("features", [])
    if features:
        total_w = len(features) * Inches(2.5)
        x_start = center_x - total_w // 2
        for i, feat in enumerate(features):
            box_x = x_start + i * Inches(2.5)
            txBox = slide.shapes.add_textbox(box_x, Inches(5.2), Inches(2.2), Inches(0.6))
            label = feat.get("label", feat) if isinstance(feat, dict) else str(feat)
            _set_text(txBox.text_frame, label, theme["font_body"], 12, theme["text_muted"])


def _render_stat_grid(slide, data: dict, theme: dict):
    """Render stat-grid: title + grid of stat boxes with values and labels."""
    import math
    _add_bg(slide, theme["primary"])
    center_x = SLIDE_WIDTH // 2

    # Title
    if data.get("title"):
        box_w = Inches(10)
        txBox = slide.shapes.add_textbox(center_x - box_w // 2, Inches(0.5), box_w, Inches(0.8))
        _set_text(txBox.text_frame, data["title"], theme["font_heading"], 36, theme["accent"], bold=True)

    stats = data.get("stats", [])
    cols = data.get("cols", min(len(stats), 4))
    if not stats or cols == 0:
        return

    rows = math.ceil(len(stats) / cols)
    card_w = Inches(2.8)
    card_h = Inches(1.8)
    gap = Inches(0.3)
    total_w = cols * card_w + (cols - 1) * gap
    x_start = center_x - total_w // 2
    y_start = Inches(1.8)

    for idx, stat in enumerate(stats):
        row = idx // cols
        col = idx % cols
        x = x_start + col * (card_w + gap)
        y = y_start + row * (card_h + gap)

        # Stat card background
        _add_rect(slide, x, y, card_w, card_h, _hex_to_rgb("#1a2236"))

        # Value
        txBox = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.3), card_w - Inches(0.4), Inches(0.8))
        value = str(stat.get("value", ""))
        _set_text(txBox.text_frame, value, theme["font_heading"], 36, theme["accent"], bold=True)

        # Label
        txBox = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1.1), card_w - Inches(0.4), Inches(0.5))
        label = stat.get("label", "")
        _set_text(txBox.text_frame, label, theme["font_body"], 12, theme["text_muted"])


def _render_two_column(slide, data: dict, theme: dict):
    """Render two-column: title + left/right columns with card/insight blocks."""
    _add_bg(slide, theme["primary"])
    center_x = SLIDE_WIDTH // 2

    # Title
    if data.get("title"):
        box_w = Inches(12)
        txBox = slide.shapes.add_textbox(center_x - box_w // 2, Inches(0.3), box_w, Inches(0.8))
        _set_text(txBox.text_frame, data["title"], theme["font_heading"], 36, theme["accent"], bold=True,
                  alignment=PP_ALIGN.LEFT)

    col_w = Inches(5.8)
    left_x = Inches(0.5)
    right_x = Inches(7.0)
    y_start = Inches(1.4)

    def render_column(blocks, x_start):
        y = y_start
        for block in blocks:
            btype = block.get("type", "card")
            if btype == "card":
                card_h = Inches(0.8) + Inches(0.25) * len(block.get("bullets", []))
                # Card background
                _add_rect(slide, x_start, y, col_w, card_h, _hex_to_rgb("#1a2236"))
                # Card title
                if block.get("title"):
                    txBox = slide.shapes.add_textbox(
                        x_start + Inches(0.2), y + Inches(0.1), col_w - Inches(0.4), Inches(0.4))
                    _set_text(txBox.text_frame, block["title"], theme["font_subheading"], 14, theme["accent"],
                              bold=True, alignment=PP_ALIGN.LEFT)
                # Card text
                if block.get("text"):
                    txBox = slide.shapes.add_textbox(
                        x_start + Inches(0.2), y + Inches(0.45), col_w - Inches(0.4), Inches(0.35))
                    _set_text(txBox.text_frame, block["text"], theme["font_body"], 11, theme["white"],
                              alignment=PP_ALIGN.LEFT)
                    txBox.text_frame.word_wrap = True
                y += card_h + Inches(0.15)
            elif btype == "insight":
                txBox = slide.shapes.add_textbox(x_start, y, col_w, Inches(0.5))
                _set_text(txBox.text_frame, block.get("text", ""), theme["font_body"], 11, theme["text_muted"],
                          alignment=PP_ALIGN.LEFT)
                txBox.text_frame.word_wrap = True
                y += Inches(0.55)
            else:
                # Unsupported block type in native — skip with spacing
                y += Inches(0.5)

    render_column(data.get("left", []), left_x)
    render_column(data.get("right", []), right_x)


def _render_three_column(slide, data: dict, theme: dict):
    """Render three-column: title + 3 columns with card blocks."""
    _add_bg(slide, theme["primary"])
    center_x = SLIDE_WIDTH // 2

    # Title
    if data.get("title"):
        box_w = Inches(12)
        txBox = slide.shapes.add_textbox(center_x - box_w // 2, Inches(0.3), box_w, Inches(0.8))
        _set_text(txBox.text_frame, data["title"], theme["font_heading"], 36, theme["accent"], bold=True,
                  alignment=PP_ALIGN.LEFT)

    columns = data.get("columns", [[], [], []])
    num_cols = len(columns)
    if num_cols == 0:
        return

    col_w = Inches(3.6)
    gap = Inches(0.3)
    total_w = num_cols * col_w + (num_cols - 1) * gap
    x_start = center_x - total_w // 2
    y_start = Inches(1.4)

    for col_idx, col_blocks in enumerate(columns):
        x = x_start + col_idx * (col_w + gap)
        y = y_start
        if not isinstance(col_blocks, list):
            continue
        for block in col_blocks:
            btype = block.get("type", "card")
            if btype == "card":
                card_h = Inches(0.8) + Inches(0.25) * len(block.get("bullets", []))
                _add_rect(slide, x, y, col_w, card_h, _hex_to_rgb("#1a2236"))
                if block.get("title"):
                    txBox = slide.shapes.add_textbox(
                        x + Inches(0.15), y + Inches(0.1), col_w - Inches(0.3), Inches(0.4))
                    _set_text(txBox.text_frame, block["title"], theme["font_subheading"], 13, theme["accent"],
                              bold=True, alignment=PP_ALIGN.LEFT)
                if block.get("text"):
                    txBox = slide.shapes.add_textbox(
                        x + Inches(0.15), y + Inches(0.45), col_w - Inches(0.3), Inches(0.3))
                    _set_text(txBox.text_frame, block["text"], theme["font_body"], 10, theme["white"],
                              alignment=PP_ALIGN.LEFT)
                    txBox.text_frame.word_wrap = True
                y += card_h + Inches(0.15)
            else:
                y += Inches(0.5)


def _render_data_table(slide, data: dict, theme: dict):
    """Render data-table: section titles + table rows + optional insight text."""
    _add_bg(slide, theme["primary"])
    center_x = SLIDE_WIDTH // 2

    # Slide title (from header partial — data.title)
    if data.get("title"):
        box_w = Inches(12)
        txBox = slide.shapes.add_textbox(center_x - box_w // 2, Inches(0.25), box_w, Inches(0.7))
        _set_text(txBox.text_frame, data["title"], theme["font_heading"], 32, theme["accent"],
                  bold=True, alignment=PP_ALIGN.LEFT)

    sections = data.get("sections", [])
    if not sections:
        return

    # Determine column layout: support up to 2 sections side by side
    num_cols = min(len(sections), 2)
    col_w = Inches(5.9) if num_cols == 2 else Inches(12)
    col_gap = Inches(0.5)
    col_xs = []
    if num_cols == 2:
        total_w = num_cols * col_w + (num_cols - 1) * col_gap
        x_start = center_x - total_w // 2
        col_xs = [x_start + i * (col_w + col_gap) for i in range(num_cols)]
    else:
        col_xs = [center_x - col_w // 2]

    y_global_start = Inches(1.15)

    for sec_idx, section in enumerate(sections[:2]):
        col_x = col_xs[sec_idx]
        y = y_global_start

        # Section title
        if section.get("section_title"):
            txBox = slide.shapes.add_textbox(col_x, y, col_w, Inches(0.45))
            _set_text(txBox.text_frame, section["section_title"], theme["font_subheading"],
                      13, theme["accent"], bold=True, alignment=PP_ALIGN.LEFT)
            y += Inches(0.5)

        # Table
        table_data = section.get("table")
        if table_data:
            headers = table_data.get("headers", [])
            rows = table_data.get("rows", [])
            winner_rows = set(table_data.get("winner_rows", []))
            total_row = table_data.get("total_row")

            all_rows = [headers] + rows
            if total_row:
                all_rows.append(total_row)

            if all_rows and headers:
                num_table_cols = len(headers)
                cell_w = col_w / num_table_cols
                row_h = Inches(0.3)

                for row_idx, row in enumerate(all_rows):
                    is_header = row_idx == 0
                    is_winner = (row_idx - 1) in winner_rows  # rows are 1-indexed after header
                    is_total = total_row and row_idx == len(all_rows) - 1

                    # Row background
                    if is_header:
                        bg_color = _hex_to_rgb("#1a2236")
                    elif is_winner:
                        bg_color = _hex_to_rgb("#1e3a1e")  # muted green tint
                    elif is_total:
                        bg_color = _hex_to_rgb("#1a2236")
                    else:
                        bg_color = _hex_to_rgb("#111827") if row_idx % 2 == 1 else None

                    if bg_color:
                        _add_rect(slide, col_x, y, col_w, row_h, bg_color)

                    for col_idx, cell in enumerate(row[:num_table_cols]):
                        cell_x = col_x + col_idx * cell_w
                        txBox = slide.shapes.add_textbox(
                            cell_x + Inches(0.05), y + Inches(0.03),
                            cell_w - Inches(0.1), row_h - Inches(0.05)
                        )
                        cell_color = theme["accent"] if is_header else (
                            theme["accent"] if is_total else theme["white"]
                        )
                        cell_align = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.RIGHT
                        _set_text(txBox.text_frame, str(cell), theme["font_body"],
                                  10, cell_color, bold=is_header, alignment=cell_align)
                    y += row_h

                if table_data.get("footnote"):
                    y += Inches(0.05)
                    txBox = slide.shapes.add_textbox(col_x, y, col_w, Inches(0.3))
                    _set_text(txBox.text_frame, table_data["footnote"], theme["font_body"],
                              9, theme["text_muted"], alignment=PP_ALIGN.LEFT)
                    y += Inches(0.35)

        # Insight
        insight = section.get("insight")
        if insight and insight.get("text"):
            y += Inches(0.1)
            txBox = slide.shapes.add_textbox(col_x, y, col_w, Inches(0.45))
            _set_text(txBox.text_frame, f"  {insight['text']}", theme["font_body"],
                      10, theme["text_muted"], alignment=PP_ALIGN.LEFT)
            txBox.text_frame.word_wrap = True


def _render_image(slide, data: dict, theme: dict):
    """Render image layout: embed picture with title and caption as native text.

    Supports full-bleed (position='full') and split (position='split') modes.
    If image path is a local file, embeds it natively via add_picture().
    If image is a remote URL or unavailable, renders a colored placeholder rectangle
    with title and caption text overlaid (no network dependency at export time).
    """
    position = data.get("position", "full")
    image_src = data.get("image", "")
    title = data.get("title", "")
    caption = data.get("caption", "")

    # Try to load the image as a local file
    image_bytes = None
    if image_src and not image_src.startswith("http"):
        try:
            img_path = Path(image_src)
            if img_path.exists():
                image_bytes = img_path.read_bytes()
        except Exception:
            pass

    if position == "split":
        side = data.get("side", "left")
        img_w = SLIDE_WIDTH // 2
        text_w = SLIDE_WIDTH // 2 - Inches(0.5)
        img_x = 0 if side == "left" else SLIDE_WIDTH // 2
        text_x = SLIDE_WIDTH // 2 + Inches(0.25) if side == "left" else Inches(0.25)

        # Image panel
        _add_bg(slide, theme["primary"])
        if image_bytes:
            try:
                slide.shapes.add_picture(
                    io.BytesIO(image_bytes), img_x, Emu(0), img_w, SLIDE_HEIGHT
                )
            except Exception:
                _add_rect(slide, img_x, Emu(0), img_w, SLIDE_HEIGHT, _hex_to_rgb("#1a2236"))
        else:
            _add_rect(slide, img_x, Emu(0), img_w, SLIDE_HEIGHT, _hex_to_rgb("#1a2236"))

        # Text panel
        y_text = Inches(2.0)
        if title:
            txBox = slide.shapes.add_textbox(text_x, y_text, text_w, Inches(1.5))
            _set_text(txBox.text_frame, title, theme["font_heading"], 36, theme["accent"],
                      bold=True, alignment=PP_ALIGN.LEFT)
            txBox.text_frame.word_wrap = True
            y_text += Inches(1.6)
        if caption:
            txBox = slide.shapes.add_textbox(text_x, y_text, text_w, Inches(1.0))
            _set_text(txBox.text_frame, caption, theme["font_body"], 14, theme["text_muted"],
                      alignment=PP_ALIGN.LEFT)
            txBox.text_frame.word_wrap = True
    else:
        # Full-bleed mode
        _add_bg(slide, theme["primary"])
        if image_bytes:
            try:
                slide.shapes.add_picture(
                    io.BytesIO(image_bytes), Emu(0), Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT
                )
            except Exception:
                _add_rect(slide, Emu(0), Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT, _hex_to_rgb("#1a2236"))
        else:
            # Placeholder: dark rectangle filling the slide
            _add_rect(slide, Emu(0), Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT, _hex_to_rgb("#1a2236"))

        # Overlay: title and caption in lower third
        if title or caption:
            # Semi-dark overlay strip at bottom
            overlay_h = Inches(1.8)
            _add_rect(slide, Emu(0), SLIDE_HEIGHT - overlay_h, SLIDE_WIDTH, overlay_h,
                      _hex_to_rgb("#0d1117"))
            y_overlay = SLIDE_HEIGHT - overlay_h + Inches(0.25)
            if title:
                txBox = slide.shapes.add_textbox(Inches(0.8), y_overlay, SLIDE_WIDTH - Inches(1.6), Inches(0.9))
                _set_text(txBox.text_frame, title, theme["font_heading"], 32, theme["white"],
                          bold=True, alignment=PP_ALIGN.LEFT)
                txBox.text_frame.word_wrap = True
                y_overlay += Inches(0.9)
            if caption:
                txBox = slide.shapes.add_textbox(Inches(0.8), y_overlay, SLIDE_WIDTH - Inches(1.6), Inches(0.5))
                _set_text(txBox.text_frame, caption, theme["font_body"], 13, theme["text_muted"],
                          alignment=PP_ALIGN.LEFT)
                txBox.text_frame.word_wrap = True


def _render_timeline(slide, data: dict, theme: dict):
    """Render timeline: title + horizontal steps with title and description text boxes."""
    _add_bg(slide, theme["primary"])
    center_x = SLIDE_WIDTH // 2

    # Title (from data.title, rendered by header partial in HTML)
    if data.get("title"):
        box_w = Inches(12)
        txBox = slide.shapes.add_textbox(center_x - box_w // 2, Inches(0.3), box_w, Inches(0.7))
        _set_text(txBox.text_frame, data["title"], theme["font_heading"], 36, theme["accent"],
                  bold=True, alignment=PP_ALIGN.LEFT)

    steps = data.get("steps", [])
    if not steps:
        return

    num_steps = len(steps)
    # Layout: steps distributed across slide width
    step_area_w = SLIDE_WIDTH - Inches(1.0)
    step_w = step_area_w / num_steps
    x_start = Inches(0.5)
    dot_y = Inches(2.8)  # vertical position of the step dots
    line_y = dot_y + Inches(0.18)  # center of connecting line
    dot_size = Inches(0.36)

    # Horizontal connecting line (drawn first, behind dots)
    line_x_start = x_start + step_w * 0.5
    line_x_end = x_start + step_w * (num_steps - 0.5)
    if num_steps > 1:
        line_w = line_x_end - line_x_start
        _add_rect(slide, line_x_start, line_y - Inches(0.02), line_w, Inches(0.04), theme["accent"])

    for i, step in enumerate(steps):
        step_x_center = x_start + step_w * i + step_w / 2
        step_x = step_x_center - dot_size / 2

        # Step dot (circle approximated as square with accent color)
        _add_rect(slide, int(step_x), int(dot_y), int(dot_size), int(dot_size), theme["accent"])

        # Step number in dot
        txBox = slide.shapes.add_textbox(
            int(step_x), int(dot_y), int(dot_size), int(dot_size)
        )
        _set_text(txBox.text_frame, str(i + 1), theme["font_heading"], 11,
                  theme["primary"], bold=True, alignment=PP_ALIGN.CENTER)

        # Step title (below dot)
        title_y = dot_y + dot_size + Inches(0.2)
        title_w = step_w - Inches(0.1)
        title_x = step_x_center - title_w / 2
        txBox = slide.shapes.add_textbox(int(title_x), int(title_y), int(title_w), Inches(0.5))
        _set_text(txBox.text_frame, step.get("title", ""), theme["font_subheading"], 12,
                  theme["white"], bold=True, alignment=PP_ALIGN.CENTER)
        txBox.text_frame.word_wrap = True

        # Step description (below title)
        desc_y = title_y + Inches(0.55)
        txBox = slide.shapes.add_textbox(int(title_x), int(desc_y), int(title_w), Inches(1.2))
        _set_text(txBox.text_frame, step.get("description", ""), theme["font_body"], 10,
                  theme["text_muted"], alignment=PP_ALIGN.CENTER)
        txBox.text_frame.word_wrap = True


def _pygments_token_color(token_type) -> RGBColor | None:
    """Return a GitHub-dark-inspired color for a Pygments token type, or None.

    Walks the token's type hierarchy (Name.Function → Name → Token) so that
    sub-types inherit the closest registered color. `None` means "use default".
    """
    try:
        from pygments.token import Token
    except ImportError:
        return None

    # Colors tuned against the #1a2236 code-block background used below.
    table = {
        Token.Comment: RGBColor(0x8B, 0x94, 0x9E),
        Token.Keyword: RGBColor(0xFF, 0x7B, 0x72),
        Token.Name.Function: RGBColor(0xD2, 0xA8, 0xFF),
        Token.Name.Class: RGBColor(0xFF, 0xA6, 0x57),
        Token.Name.Builtin: RGBColor(0x79, 0xC0, 0xFF),
        Token.Name.Decorator: RGBColor(0xD2, 0xA8, 0xFF),
        Token.String: RGBColor(0xA5, 0xD6, 0xFF),
        Token.Number: RGBColor(0x79, 0xC0, 0xFF),
        Token.Operator: RGBColor(0xFF, 0x7B, 0x72),
        Token.Punctuation: RGBColor(0xE6, 0xED, 0xF3),
    }
    node = token_type
    while node is not None:
        if node in table:
            return table[node]
        node = node.parent
    return None


def _lex_code(code: str, language: str):
    """Return a list of (token_type, text) tuples. None if Pygments unavailable.

    Unknown languages fall back to `TextLexer`, which emits one Token.Text
    span — the caller will just render one uncolored run.
    """
    try:
        from pygments import lex
        from pygments.lexers import get_lexer_by_name
        from pygments.lexers.special import TextLexer
        from pygments.util import ClassNotFound
    except ImportError:
        return None
    try:
        lexer = get_lexer_by_name(language) if language else TextLexer()
    except ClassNotFound:
        lexer = TextLexer()
    return list(lex(code, lexer))


def _render_code(slide, data: dict, theme: dict):
    """Render code: title + single mono text box with source, optional caption.

    Uses Pygments for per-token colors when available; falls back to a single
    uncolored run otherwise (Pygments is an optional dependency — see
    docs/PLAN.md T2.2, AH-G2).
    """
    _add_bg(slide, theme["primary"])
    center_x = SLIDE_WIDTH // 2

    # Title (rendered as native text so editors can retitle the slide)
    if data.get("title"):
        box_w = Inches(12)
        txBox = slide.shapes.add_textbox(
            center_x - box_w // 2, Inches(0.3), box_w, Inches(0.7)
        )
        _set_text(
            txBox.text_frame, data["title"], theme["font_heading"],
            28, theme["accent"], bold=True, alignment=PP_ALIGN.LEFT,
        )

    has_caption = bool(data.get("caption"))
    code_x = Inches(0.5)
    code_y = Inches(1.15)
    code_w = SLIDE_WIDTH - Inches(1.0)
    code_h = Inches(5.3) if has_caption else Inches(5.9)

    # Dark code-block background (distinct from theme primary so the box reads
    # as a terminal/editor pane on any preset)
    _add_rect(slide, code_x, code_y, code_w, code_h, _hex_to_rgb("#1a2236"))

    # Code text box — single editable shape with per-token runs
    pad = Inches(0.25)
    txBox = slide.shapes.add_textbox(
        code_x + pad, code_y + pad, code_w - pad * 2, code_h - pad * 2
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    _populate_code_frame(
        tf,
        code=data.get("code", ""),
        language=data.get("language", ""),
        font_mono=theme["font_mono"],
    )

    # Caption (optional)
    if has_caption:
        cap_y = code_y + code_h + Inches(0.12)
        txBox = slide.shapes.add_textbox(code_x, cap_y, code_w, Inches(0.4))
        _set_text(
            txBox.text_frame, data["caption"], theme["font_body"],
            12, theme["text_muted"], alignment=PP_ALIGN.LEFT,
        )


def _populate_code_frame(text_frame, code: str, language: str, font_mono: str):
    """Fill a text_frame with the code, one paragraph per line, colored runs.

    Falls back to a single uncolored run if Pygments is unavailable or the
    language lexer can't be resolved.
    """
    default_color = RGBColor(0xE6, 0xED, 0xF3)
    font_size = Pt(14)

    text_frame.clear()
    # Remove the auto-created empty paragraph; we'll add paragraphs per line
    # below to preserve explicit blank lines in the source.

    tokens = _lex_code(code, language)
    if tokens is None:
        # No Pygments — single run, newlines handled as separate paragraphs
        lines = code.splitlines() or [""]
        for i, line in enumerate(lines):
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = line
            run.font.name = font_mono
            run.font.size = font_size
            run.font.color.rgb = default_color
        return

    # Split token stream into per-line (token_type, text) segments so each
    # source line becomes its own paragraph. Preserves blank lines.
    line_segments: list[list[tuple]] = [[]]
    for token_type, text in tokens:
        if not text:
            continue
        parts = text.split("\n")
        for idx, part in enumerate(parts):
            if part:
                line_segments[-1].append((token_type, part))
            if idx < len(parts) - 1:
                line_segments.append([])

    for i, segments in enumerate(line_segments):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if not segments:
            # Blank line — add a zero-width run so the paragraph occupies vertical space
            run = p.add_run()
            run.text = ""
            run.font.name = font_mono
            run.font.size = font_size
            continue
        for token_type, text in segments:
            run = p.add_run()
            run.text = text
            run.font.name = font_mono
            run.font.size = font_size
            color = _pygments_token_color(token_type)
            run.font.color.rgb = color if color is not None else default_color


# ── Chart type mapping (T2.4, D5) ────────────────────────────────
#
# Plotly's `type: bar` renders vertical by default (a "column chart" in
# PowerPoint's terminology). For visual parity we map "bar" to
# COLUMN_CLUSTERED even though the plan text says BAR_CLUSTERED —
# BAR_CLUSTERED is horizontal in Excel/PowerPoint and would surprise
# existing decks that already use `chart_type: bar`. "bar-horizontal"
# is exposed as an explicit opt-in. Bedrock plan mapping (line, pie,
# donut/doughnut, scatter, area) is unchanged.

def _xl_chart_type(name: str):
    """Resolve a YAML chart_type string to a python-pptx XL_CHART_TYPE or None."""
    from pptx.enum.chart import XL_CHART_TYPE
    mapping = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,        # Plotly default: vertical
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "bar-horizontal": XL_CHART_TYPE.BAR_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
        "donut": XL_CHART_TYPE.DOUGHNUT,
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
        "scatter": XL_CHART_TYPE.XY_SCATTER,
        "area": XL_CHART_TYPE.AREA,
    }
    return mapping.get((name or "").lower())


def _apply_chart_theme(chart, theme: dict) -> None:
    """Color each series with accent → secondary_accent → accent → …

    Uses both fill (for bar/column/area/pie/donut) and line (for line/scatter)
    properties — whichever API the series type doesn't support raises, so we
    swallow the exception and move on rather than type-sniffing.
    """
    palette = [theme["accent"]]
    if theme.get("secondary_accent"):
        palette.append(theme["secondary_accent"])

    try:
        series_list = list(chart.series)
    except Exception:
        return
    for i, series in enumerate(series_list):
        color = palette[i % len(palette)]
        try:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = color
        except Exception:
            pass
        try:
            series.format.line.color.rgb = color
        except Exception:
            pass


def _render_chart(slide, data: dict, theme: dict):
    """Render a native PowerPoint chart (double-clickable, editable data table).

    Consumes the same YAML shape as the HTML chart template:
      - ``chart_type``: bar | column | line | pie | donut | scatter | area
      - ``labels``: list of category names (x-axis)
      - ``values``: list of numbers (single-series) — OR —
      - ``series``: list of ``{name, values}`` dicts (multi-series)
      - ``title``, ``subtitle``: optional header text

    Unknown types or empty data render a captioned placeholder text box so
    the slide still has an editable shape; T2.8's --strict tracking will
    surface this as a fallback event.
    """
    _add_bg(slide, theme["primary"])
    center_x = SLIDE_WIDTH // 2

    # Title
    top = Inches(0.3)
    if data.get("title"):
        box_w = Inches(12)
        txBox = slide.shapes.add_textbox(center_x - box_w // 2, top, box_w, Inches(0.7))
        _set_text(
            txBox.text_frame, data["title"], theme["font_heading"],
            32, theme["accent"], bold=True, alignment=PP_ALIGN.LEFT,
        )
        top = Inches(1.0)

    # Subtitle (optional)
    if data.get("subtitle"):
        box_w = Inches(12)
        txBox = slide.shapes.add_textbox(center_x - box_w // 2, top, box_w, Inches(0.5))
        _set_text(
            txBox.text_frame, data["subtitle"], theme["font_body"],
            14, theme["text_muted"], alignment=PP_ALIGN.LEFT,
        )
        top += Inches(0.55)

    chart_type_name = (data.get("chart_type") or "bar").lower()
    xl_type = _xl_chart_type(chart_type_name)

    labels = data.get("labels") or []
    values = data.get("values") or []
    series_defs = data.get("series") or []

    chart_x = Inches(0.5)
    chart_y = top + Inches(0.2)
    chart_w = SLIDE_WIDTH - Inches(1.0)
    chart_h = SLIDE_HEIGHT - chart_y - Inches(0.5)

    if xl_type is None or (not labels and not series_defs and not values):
        # Can't build a real chart — leave an editable caption. The coverage
        # gate (T2.1) still passes on the text frame; T2.8 will treat this
        # as a fallback and fail --strict.
        txBox = slide.shapes.add_textbox(chart_x, chart_y, chart_w, Inches(1.2))
        _set_text(
            txBox.text_frame,
            f"[Chart: {chart_type_name or 'unknown'} — unrenderable in native PPTX]",
            theme["font_body"], 16, theme["text_muted"], alignment=PP_ALIGN.CENTER,
        )
        return

    from pptx.chart.data import CategoryChartData, XyChartData
    from pptx.enum.chart import XL_CHART_TYPE

    if xl_type == XL_CHART_TYPE.XY_SCATTER:
        chart_data = XyChartData()
        if series_defs:
            for s in series_defs:
                name = s.get("name", "Series")
                pts = s.get("xy")
                if pts is None:
                    pts = list(zip(s.get("x", []), s.get("y", [])))
                series = chart_data.add_series(name)
                for pt in pts:
                    try:
                        series.add_data_point(float(pt[0]), float(pt[1]))
                    except (TypeError, ValueError, IndexError):
                        continue
        elif labels and values:
            series = chart_data.add_series(data.get("name", "Series"))
            for x, y in zip(labels, values):
                try:
                    series.add_data_point(float(x), float(y))
                except (TypeError, ValueError):
                    continue
    else:
        chart_data = CategoryChartData()
        chart_data.categories = [str(x) for x in labels] or [""]
        if series_defs:
            for s in series_defs:
                chart_data.add_series(s.get("name", "Series"), s.get("values", []))
        else:
            chart_data.add_series(data.get("name", "Series"), values)

    chart_shape = slide.shapes.add_chart(
        xl_type, chart_x, chart_y, chart_w, chart_h, chart_data
    )
    _apply_chart_theme(chart_shape.chart, theme)


def _render_map(slide, data: dict, theme: dict, *, slide_file=None, pw_context=None):
    """Render the map layout: rasterized map image on the left, editable legend
    with one text box per marker on the right.

    Leaflet is a live JS widget rendered in a browser — we can't round-trip it
    to native PowerPoint shapes. The plan (T2.5) accepts this: embed the
    Playwright-rendered screenshot for the visual, then surface the markers
    (``data.markers[]`` with lat/lng/label) as a right-edge legend so each
    annotation stays editable.
    """
    _add_bg(slide, theme["primary"])

    markers = data.get("markers") or []

    # Map image on the left ~75% of the slide; legend gets the right ~25%
    # when there are markers, otherwise the image takes the full width.
    has_legend = bool(markers)
    map_w = Inches(9.5) if has_legend else SLIDE_WIDTH
    legend_x = map_w + Inches(0.1) if has_legend else None
    legend_w = SLIDE_WIDTH - (legend_x or Emu(0)) - Inches(0.2) if has_legend else Emu(0)

    # Title (rendered before the image so it shows over the dark panel, not
    # over the map itself — the slide_file render already includes it in
    # the raster but PPTX viewers can't edit pixels).
    title_h = Inches(0.5)
    if data.get("title"):
        txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), map_w - Inches(0.6), title_h)
        _set_text(
            txBox.text_frame, data["title"], theme["font_heading"],
            24, theme["accent"], bold=True, alignment=PP_ALIGN.LEFT,
        )

    # Rasterize the map slide and embed. If Playwright isn't available, skip
    # the image and let the legend carry the slide.
    if slide_file is not None and slide_file.exists():
        try:
            img_bytes = _rasterize_slide(slide_file, pw_context)
        except Exception:
            img_bytes = None
        if img_bytes:
            # Crop to map area: embed at full map_w width, scale height to
            # keep aspect ratio. The HTML-rendered map fills the whole slide,
            # so scaling to slide_height is OK here too.
            slide.shapes.add_picture(
                io.BytesIO(img_bytes),
                Emu(0), Emu(0), map_w, SLIDE_HEIGHT,
            )

    # Editable legend on the right
    if has_legend:
        legend_y = Inches(0.15)
        # Panel background so the legend reads cleanly over the theme primary
        _add_rect(slide, legend_x, Emu(0), legend_w + Inches(0.2), SLIDE_HEIGHT,
                  _hex_to_rgb("#1a2236"))
        # Legend header
        txBox = slide.shapes.add_textbox(legend_x + Inches(0.1), legend_y,
                                          legend_w, Inches(0.5))
        _set_text(
            txBox.text_frame, "Markers", theme["font_subheading"],
            16, theme["accent"], bold=True, alignment=PP_ALIGN.LEFT,
        )
        legend_y += Inches(0.55)
        entry_h = Inches(0.55)
        for m in markers:
            label = m.get("label") or f"{m.get('lat', '?')}, {m.get('lng', '?')}"
            if legend_y + entry_h > SLIDE_HEIGHT - Inches(0.2):
                break
            txBox = slide.shapes.add_textbox(legend_x + Inches(0.1), legend_y,
                                              legend_w, entry_h)
            tf = txBox.text_frame
            tf.word_wrap = True
            _set_text(
                tf, f"• {label}", theme["font_body"],
                12, theme["white"], alignment=PP_ALIGN.LEFT,
            )
            legend_y += entry_h + Inches(0.05)


def _rasterize_slide(slide_file, pw_context):
    """Screenshot a slide HTML file. Uses a shared Playwright context when
    provided; otherwise spawns its own. Returns PNG bytes or None.
    """
    if pw_context is not None:
        page = pw_context.new_page()
        try:
            page.goto(f"file://{slide_file.resolve()}")
            page.wait_for_load_state("networkidle")
            try:
                page.wait_for_selector("[data-pf-ready]", timeout=10000)
            except Exception:
                pass
            return page.screenshot(full_page=False)
        finally:
            page.close()

    try:
        from playwright.sync_api import sync_playwright
    except ImportError:
        return None
    with sync_playwright() as p:
        browser = p.chromium.launch()
        try:
            page = browser.new_page(viewport={"width": 1280, "height": 720})
            page.goto(f"file://{slide_file.resolve()}")
            page.wait_for_load_state("networkidle")
            try:
                page.wait_for_selector("[data-pf-ready]", timeout=10000)
            except Exception:
                pass
            return page.screenshot(full_page=False)
        finally:
            browser.close()


def _render_toc(slide, data: dict, theme: dict, *, prs=None, slides_cfg=None, slide_index: int = 0):
    """Render table of contents: one editable text box per entry, each with
    a NAMED_SLIDE click_action hyperlink to the matching section slide.

    Entry → target slide resolution:
      1. If the entry has an explicit ``slide`` field (1-based per YAML
         convention), use it.
      2. Otherwise, map the nth TOC entry to the nth ``section``-layout slide
         in ``slides_cfg`` — this matches the behavior of
         ``PresentationBuilder._generate_toc``.
    """
    _add_bg(slide, theme["primary"])
    center_x = SLIDE_WIDTH // 2

    # Slide title (via header partial in HTML; rendered as native text here)
    if data.get("title"):
        box_w = Inches(12)
        txBox = slide.shapes.add_textbox(
            center_x - box_w // 2, Inches(0.3), box_w, Inches(0.8)
        )
        _set_text(
            txBox.text_frame, data["title"], theme["font_heading"],
            36, theme["accent"], bold=True, alignment=PP_ALIGN.LEFT,
        )

    items = data.get("items") or []
    if not items:
        return

    # Sequential section-slide indices for fallback mapping
    section_indices: list[int] = []
    if slides_cfg:
        section_indices = [
            i for i, sc in enumerate(slides_cfg) if sc.get("layout") == "section"
        ]

    # Vertical layout of TOC entries
    entry_h = Inches(0.55)
    entry_gap = Inches(0.1)
    total_h = len(items) * entry_h + max(0, len(items) - 1) * entry_gap
    y_start = max(Inches(1.4), (SLIDE_HEIGHT - total_h) // 2)
    entry_w = Inches(10)
    x_start = center_x - entry_w // 2

    for idx, item in enumerate(items):
        y = y_start + idx * (entry_h + entry_gap)
        txBox = slide.shapes.add_textbox(x_start, y, entry_w, entry_h)
        tf = txBox.text_frame
        tf.word_wrap = True

        number = item.get("number")
        if number is None:
            number_str = f"{idx + 1:02d}"
        elif isinstance(number, int):
            number_str = f"{number:02d}"
        else:
            number_str = str(number)
        title = item.get("title", "")
        text = f"{number_str} — {title}" if title else number_str
        _set_text(
            tf, text, theme["font_subheading"],
            20, theme["white"], bold=False, alignment=PP_ALIGN.LEFT,
        )

        # Resolve hyperlink target
        if prs is None:
            continue
        target_idx: int | None = None
        raw = item.get("slide")
        if isinstance(raw, int) and raw >= 1:
            target_idx = raw - 1
        elif idx < len(section_indices):
            target_idx = section_indices[idx]
        if target_idx is not None and 0 <= target_idx < len(prs.slides):
            txBox.click_action.target_slide = prs.slides[target_idx]


# ── Layout dispatch ──────────────────────────────────────────────

NATIVE_RENDERERS = {
    "section": _render_section,
    "quote": _render_quote,
    "closing": _render_closing,
    "title": _render_title,
    "stat-grid": _render_stat_grid,
    "two-column": _render_two_column,
    "three-column": _render_three_column,
    "data-table": _render_data_table,
    "image": _render_image,
    "timeline": _render_timeline,
    "code": _render_code,
    "toc": _render_toc,
    "chart": _render_chart,
    "map": _render_map,
}


def _render_image_fallback(slide, slide_file: Path, context=None):
    """Fall back to screenshot for complex layouts (requires Playwright).

    Args:
        slide: python-pptx slide object to add the screenshot to.
        slide_file: Path to the slide HTML file.
        context: Optional shared Playwright browser context. If provided,
            reuses it instead of spawning a new browser per slide.
    """
    if context:
        page = context.new_page()
        page.goto(f"file://{slide_file.resolve()}")
        page.wait_for_load_state("networkidle")
        try:
            page.wait_for_selector("[data-pf-ready]", timeout=10000)
        except Exception:
            pass  # Graceful fallback if sentinel missing
        png_bytes = page.screenshot(full_page=False)
        page.close()
    else:
        # Legacy path — spawns own browser (no shared context available)
        try:
            from playwright.sync_api import sync_playwright
        except ImportError:
            return  # Skip if Playwright unavailable

        with sync_playwright() as p:
            browser = p.chromium.launch()
            page = browser.new_page(viewport={"width": 1280, "height": 720})
            page.goto(f"file://{slide_file.resolve()}")
            page.wait_for_load_state("networkidle")
            try:
                page.wait_for_selector("[data-pf-ready]", timeout=10000)
            except Exception:
                pass  # Graceful fallback if sentinel missing
            png_bytes = page.screenshot(full_page=False)
            page.close()
            browser.close()

    slide.shapes.add_picture(
        io.BytesIO(png_bytes),
        left=Emu(0), top=Emu(0),
        width=SLIDE_WIDTH, height=SLIDE_HEIGHT,
    )


def export_pptx_editable(
    config: dict,
    slides_dir: str,
    output_path: str,
):
    """Export slides to an editable .pptx file.

    Native text/shapes for supported layouts (section, quote, closing,
    title, stat-grid, two-column, three-column). Image fallback via
    Playwright for complex layouts. Uses a single shared browser context
    across all image fallbacks to avoid spawning a browser per slide.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]

    theme = _pptx_theme(config.get("theme", {}))
    slides_path = Path(slides_dir)
    slides_cfg = config.get("slides", [])

    # Create shared browser context for image fallbacks (EXPORT-04)
    pw_manager = None
    pw_browser = None
    pw_context = None
    try:
        from playwright.sync_api import sync_playwright
        pw_manager = sync_playwright().start()
        pw_browser = pw_manager.chromium.launch()
        pw_context = pw_browser.new_context(viewport={"width": 1280, "height": 720})
    except (ImportError, Exception):
        pass  # Playwright unavailable — image fallback will skip gracefully

    try:
        # Two-pass so renderers (e.g. TOC) can reference forward-indexed slides
        # via prs.slides[N] when they run — otherwise a TOC slide at position 0
        # can't link to a section slide that hasn't been added yet.
        slide_objs = [prs.slides.add_slide(blank_layout) for _ in slides_cfg]

        for i, slide_cfg in enumerate(slides_cfg):
            layout = slide_cfg.get("layout", "two-column")
            data = slide_cfg.get("data", {})
            slide = slide_objs[i]
            slide_file = slides_path / f"slide_{i + 1:02d}.html"

            renderer = NATIVE_RENDERERS.get(layout)
            if renderer:
                # Only pass context kwargs the renderer actually declares, so
                # existing renderers with signature (slide, data, theme) stay
                # untouched while TOC / map / video can opt in to extras.
                params = inspect.signature(renderer).parameters
                extras = {}
                if "prs" in params:
                    extras["prs"] = prs
                if "slides_cfg" in params:
                    extras["slides_cfg"] = slides_cfg
                if "slide_index" in params:
                    extras["slide_index"] = i
                if "slide_file" in params:
                    extras["slide_file"] = slide_file
                if "pw_context" in params:
                    extras["pw_context"] = pw_context
                renderer(slide, data, theme, **extras)
            else:
                if slide_file.exists():
                    _render_image_fallback(slide, slide_file, context=pw_context)

            # Speaker notes
            if slide_cfg.get("notes"):
                slide.notes_slide.notes_text_frame.text = slide_cfg["notes"]
    finally:
        # Always clean up shared browser context
        if pw_context:
            pw_context.close()
        if pw_browser:
            pw_browser.close()
        if pw_manager:
            pw_manager.stop()

    prs.save(output_path)
